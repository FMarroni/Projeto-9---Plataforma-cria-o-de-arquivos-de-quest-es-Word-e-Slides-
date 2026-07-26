"""Gera os templates `templates/lista.docx`, `templates/comentada.docx` e
`templates/slides.pptx` a partir da identidade visual corporativa em
`Templates_padrão/` (Estratégia Concursos): o PPTX é derivado de
`Templates_padrão/Slides Mestre.pptx` (mestre/layouts/logotipos reais,
preservados), e o DOCX usa a cor de destaque (#4231A4) e a fonte (Cambria)
extraídas dos exemplos `ep48Rg5nBQ1zIsJVGiwI (comentarios/lista).docx`.

O relatório de análise agregada NÃO é um docx — vira HTML, gerado diretamente
por `html_gen.py` (sem template).

Este script pode ser reexecutado a qualquer momento para recriar os
templates-base (ex.: se a identidade visual mudar) — o código de geração
(`docx_lista.py`/`docx_comentada.py`/`pptx_gen.py`) só depende dos tokens
`{{...}}` e dos nomes de shape/tabela citados abaixo, não do design visual em
si; PORÉM os índices de slide (CAPA_IDX etc., em `app/pptx_gen.py`) dependem
da estrutura de slides que este script produz — ver o comentário no final.

Rode com: backend/.venv/Scripts/python.exe backend/scripts/build_templates.py
"""

import os
import sys

from docx import Document
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PCm
from pptx.util import Pt as PPt

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
TEMPLATES_DIR = os.path.join(ROOT, "templates")
TEMPLATES_PADRAO_DIR = os.path.join(ROOT, "Templates_padrão")
SLIDES_MESTRE_PATH = os.path.join(TEMPLATES_PADRAO_DIR, "Slides Mestre.pptx")

sys.path.insert(0, os.path.join(ROOT, "backend"))
from app.pptx_xml_utils import duplicar_slide, remover_slide  # noqa: E402

# --- Identidade visual (Estratégia Concursos) — extraída de Templates_padrão/ ---
# Cor de destaque: run bold "(CEBRASPE / TRF 6ª Região - 2025)" / "Comentários:"
# nos .docx de exemplo. Fonte do corpo: Cambria (docDefaults -> minorFont do
# theme desses .docx). Fonte de títulos/UI no PPTX: Montserrat (run "Prof."
# no layout Capa_Disciplina do Slides Mestre.pptx).
COR_MARCA = RGBColor(0x42, 0x31, 0xA4)
COR_MARCA_PPTX = PptxRGBColor(0x42, 0x31, 0xA4)
FONTE_CORPO_DOCX = "Cambria"
FONTE_UI_PPTX = "Montserrat"

COR_HEADER_BG = "D9D2E9"  # tom claro derivado da cor de marca, p/ shading de célula
COR_DESTAQUE_BG = "FCE8B2"  # amarelo suave p/ Curva ABC (hex sem '#', p/ shading XML)


def _set_cell_shading(cell, hex_color: str) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    shd = tcPr.makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:fill"): hex_color})
    tcPr.append(shd)


def _set_cell_text(cell, texto: str, bold: bool = False, color: RGBColor | None = None) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(texto)
    run.font.name = FONTE_CORPO_DOCX
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_run_padrao(run, bold: bool = False, color: RGBColor | None = None, size: Pt | None = None) -> None:
    run.font.name = FONTE_CORPO_DOCX
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if size is not None:
        run.font.size = size


def _aplicar_estilos_documento(doc: Document) -> None:
    """Define a fonte de marca como padrão do documento (estilo "Normal") e a
    cor de marca nos headings — assim, texto adicionado em runtime por
    docx_lista.py/docx_comentada.py (que não define fonte/cor por run) herda
    a identidade visual automaticamente, sem precisar espalhar
    `run.font.name = ...` pelo código de geração. "Heading 1" é usado para o
    agrupamento por matéria (ver docx_lista.py/docx_comentada.py)."""
    normal = doc.styles["Normal"]
    normal.font.name = FONTE_CORPO_DOCX
    normal.font.size = Pt(11)

    for nome_estilo in ("Heading 1", "Heading 3"):
        estilo = doc.styles[nome_estilo]
        estilo.font.name = FONTE_CORPO_DOCX
        estilo.font.color.rgb = COR_MARCA


def build_lista_template() -> str:
    """Lista de questões: só enunciado + alternativas + imagens; gabarito
    consolidado ao final (como o bloco "Gabarito" do PDF original do TEC)."""
    doc = Document()
    _aplicar_estilos_documento(doc)

    _set_run_padrao(doc.add_paragraph().add_run("{{LISTA_QUESTOES}}"))

    gabarito_heading = doc.add_heading("", level=1)
    _set_run_padrao(gabarito_heading.add_run("Gabarito"), bold=True, color=COR_MARCA)
    _set_run_padrao(doc.add_paragraph().add_run("{{GABARITO_CONSOLIDADO}}"))

    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    caminho = os.path.join(TEMPLATES_DIR, "lista.docx")
    doc.save(caminho)
    return caminho


def build_comentada_template() -> str:
    """Questões comentadas: enunciado + alternativas + imagens + comentário +
    gabarito, sequencialmente, uma questão após a outra."""
    doc = Document()
    _aplicar_estilos_documento(doc)

    _set_run_padrao(doc.add_paragraph().add_run("{{QUESTOES_COMENTADAS}}"))

    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    caminho = os.path.join(TEMPLATES_DIR, "comentada.docx")
    doc.save(caminho)
    return caminho


def build_rastreabilidade_template() -> str:
    """Épico 4 — Módulo Biblioteca: trilha de auditoria (arquivo+página que
    fundamentou cada alternativa julgada em Modo Restrito/RAG) — documento
    simples, de uso interno, nunca entregue ao aluno final junto dos demais."""
    doc = Document()
    _aplicar_estilos_documento(doc)

    titulo = doc.add_heading("", level=1)
    _set_run_padrao(
        titulo.add_run("Rastreabilidade — Base de Conhecimento (RAG)"), bold=True, color=COR_MARCA
    )

    _set_run_padrao(doc.add_paragraph().add_run("{{RASTREABILIDADE}}"))

    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    caminho = os.path.join(TEMPLATES_DIR, "rastreabilidade.docx")
    doc.save(caminho)
    return caminho


# --- PPTX: derivado de Templates_padrão/Slides Mestre.pptx -------------------
#
# Estrutura original do arquivo fonte (Google Slides exportado):
#   slide 0: layout "Questao_Automatica"  (2 placeholders BODY idx=10/11 + pic)
#   slide 1: layout "Capa_Disciplina"     (placeholders BODY idx=13/14 + pic)
#   slide 2: layout "Default 2"           ("Obrigado!" + pic)
#   slide 3: layout "Default 2"           (pic, sem texto — sobra/alternativo)
#
# Decisão de design: o layout "Questao_Automatica" NÃO tem imagem de fundo no
# próprio layout (só nos placeholders herdados) — a imagem vívida (gradiente
# roxo/azul) do slide 0 original é um plano de fundo ad-hoc daquela instância
# específica, não do layout reutilizável. Como o slide de questão se repete
# muitas vezes e carrega bastante texto denso, priorizamos legibilidade:
# construímos o slide de Questão do zero A PARTIR DO LAYOUT (fundo branco,
# herda os 2 placeholders + o rodapé de marca "Ranking: ..."), em vez de
# duplicar a instância com o gradiente vívido. Cover e Encerramento (usados
# uma única vez cada) mantêm o visual mais impactante das instâncias originais.
CAPA_ORIGINAL_IDX = 1  # layout "Capa_Disciplina", plano de fundo claro (image4)
QUESTAO_LAYOUT_IDX = 0  # layout "Questao_Automatica"
ENCERRAMENTO_ORIGINAL_IDX = 2  # "Obrigado!" (Default 2, image2)
SPARE_ORIGINAL_IDX = 3  # instância não utilizada — descartada

IDX_TITULO_CAPA = 13
IDX_SUBTITULO_CAPA = 14
IDX_HEADER_QUESTAO = 10
IDX_CORPO_QUESTAO = 11


def _placeholder_por_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _definir_texto_placeholder(ph, nome_shape: str, texto: str) -> None:
    ph.name = nome_shape
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = texto
    for run in tf.paragraphs[0].runs:
        run.font.name = FONTE_UI_PPTX


def _adicionar_textbox(slide, nome: str, left_cm, top_cm, width_cm, height_cm, texto: str):
    """As instâncias de slide duplicadas de 'Capa_Disciplina' (Templates_padrão)
    não têm os placeholders idx=13/14 materializados como shapes (só a imagem
    de fundo) — apenas slides criados via `add_slide(layout)` os materializam.
    Por isso Capa/Divisor usam caixas de texto simples (mesma posição/papel
    dos placeholders originais do layout), em vez de tentar localizar um
    placeholder que não existe nessas instâncias."""
    box = slide.shapes.add_textbox(PCm(left_cm), PCm(top_cm), PCm(width_cm), PCm(height_cm))
    box.name = nome
    box.text_frame.word_wrap = True
    box.text_frame.text = texto
    for run in box.text_frame.paragraphs[0].runs:
        run.font.name = FONTE_UI_PPTX
    return box


def _adicionar_tabela_incidencia(slide, left_cm, top_cm, width_cm, height_cm):
    tabela_shape = slide.shapes.add_table(2, 3, PCm(left_cm), PCm(top_cm), PCm(width_cm), PCm(height_cm))
    tabela_shape.name = "tabela_incidencia"
    tabela = tabela_shape.table
    for idx, texto in enumerate(["Assunto", "Nº Questões", "Incidência"]):
        cell = tabela.cell(0, idx)
        cell.text = texto
        cell.fill.solid()
        cell.fill.fore_color.rgb = COR_MARCA_PPTX
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.name = FONTE_UI_PPTX
        run.font.bold = True
        run.font.color.rgb = PptxRGBColor(0xFF, 0xFF, 0xFF)
    for idx, texto in enumerate(["(assunto de exemplo)", "0", "0,00%"]):
        cell = tabela.cell(1, idx)
        cell.text = texto
        cell.text_frame.paragraphs[0].runs[0].font.name = FONTE_UI_PPTX
    return tabela_shape


def build_pptx_template() -> str:
    if not os.path.exists(SLIDES_MESTRE_PATH):
        raise FileNotFoundError(
            f"Template corporativo não encontrado: {SLIDES_MESTRE_PATH}. "
            "Coloque o arquivo em Templates_padrão/Slides Mestre.pptx."
        )

    prs = Presentation(SLIDES_MESTRE_PATH)

    # 1) descarta a instância sobressalente (Default 2 sem texto)
    remover_slide(prs, SPARE_ORIGINAL_IDX)
    # estado: [0:Questao_original, 1:Capa_original, 2:Obrigado]

    # 2) duplica a instância "Capa_Disciplina" — uma vira Capa, a outra Divisor
    duplicar_slide(prs, CAPA_ORIGINAL_IDX)  # insere em CAPA_ORIGINAL_IDX+1 = 2
    # estado: [0:Questao_original, 1:Capa_original(->CAPA), 2:Capa_copia(->DIVISOR), 3:Obrigado]

    # 3) cria o slide de Questão do zero a partir do layout (fundo limpo, ver nota acima)
    layout_questao = prs.slide_masters[0].slide_layouts[QUESTAO_LAYOUT_IDX]
    prs.slides.add_slide(layout_questao)  # sempre anexado no fim -> índice 4
    # estado: [0:Questao_original, 1:Capa(CAPA), 2:Capa_copia(DIVISOR), 3:Obrigado, 4:Questao_novo]

    # 4) reposiciona o slide de Questão recém-criado para ficar antes do Encerramento
    xml_slides = prs.slides._sldIdLst
    slides_xml = list(xml_slides)
    sld_questao_novo = slides_xml[4]
    xml_slides.remove(sld_questao_novo)
    xml_slides.insert(3, sld_questao_novo)
    # estado: [0:Questao_original, 1:Capa(CAPA), 2:Capa_copia(DIVISOR), 3:Questao_novo(QUESTAO), 4:Obrigado(ENCERRAMENTO)]

    # 5) descarta a instância original do slide de Questão (já não usamos o fundo vívido dela)
    remover_slide(prs, 0)
    # estado final: [0:CAPA, 1:DIVISOR, 2:QUESTAO, 3:ENCERRAMENTO] — bate com as
    # constantes CAPA_IDX/DIVISOR_BASE_IDX/QUESTAO_BASE_IDX/ENCERRAMENTO_IDX = 0,1,2,3
    # já existentes em app/pptx_gen.py — não foi necessário alterá-las.

    slide_capa = prs.slides[0]
    slide_divisor = prs.slides[1]
    slide_questao = prs.slides[2]
    slide_encerramento = prs.slides[3]

    # --- Capa: mesma posição do placeholder original (idx=13/14 do layout
    # Capa_Disciplina), mas como caixa de texto simples (ver _adicionar_textbox) ---
    ph_titulo = _adicionar_textbox(slide_capa, "ph_titulo", 1.91, 3.48, 22.5, 3.5, "{{NOME_CONCURSO}}")
    for run in ph_titulo.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = PPt(36)
    ph_titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    ph_subtitulo = _adicionar_textbox(
        slide_capa,
        "ph_subtitulo",
        1.91,
        9.6,
        22.5,
        3.5,
        "{{NOME_CARGO}} — {{NOME_BANCA}}\n"
        "{{ESCOLARIDADE}} • {{ANOS_ANALISE}} • {{TOTAL_QUESTOES}} questões analisadas\n"
        "{{TEXTO_EDITAL}}",
    )
    for para in ph_subtitulo.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = PPt(16)

    # --- Divisor de disciplina: mesma instância/plano de fundo da Capa, mas
    # com os textos comprimidos no topo para abrir espaço para a tabela ---
    ph_disciplina = _adicionar_textbox(slide_divisor, "ph_disciplina", 1, 0.6, 23.4, 2.1, "{{NOME_DISCIPLINA}}")
    for run in ph_disciplina.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = PPt(28)

    ph_curva_abc = _adicionar_textbox(
        slide_divisor,
        "ph_curva_abc",
        1,
        2.7,
        23.4,
        1.3,
        "Primeiros assuntos (~50% de incidência): {{PRIMEIROS_ASSUNTOS}} ({{porcentagem}})",
    )
    for run in ph_curva_abc.text_frame.paragraphs[0].runs:
        run.font.italic = True
        run.font.size = PPt(14)

    _adicionar_tabela_incidencia(slide_divisor, left_cm=1, top_cm=4.2, width_cm=23.4, height_cm=9.6)

    # --- Questão: herdada do layout, já vem com os 2 placeholders certos ---
    ph_header = _placeholder_por_idx(slide_questao, IDX_HEADER_QUESTAO)
    _definir_texto_placeholder(ph_header, "ph_header_questao", "{{cabecalho da questão}}")
    for run in ph_header.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = PPt(18)

    ph_corpo = _placeholder_por_idx(slide_questao, IDX_CORPO_QUESTAO)
    _definir_texto_placeholder(ph_corpo, "ph_corpo_questao", "{{corpo da questão}}")
    ph_corpo.text_frame.word_wrap = True

    # --- Encerramento: mantém "Obrigado!" original, só injeta o token do concurso ---
    caixa_obrigado = next(
        (s for s in slide_encerramento.shapes if s.has_text_frame and s.text_frame.text.strip()), None
    )
    if caixa_obrigado is not None:
        caixa_obrigado.name = "ph_encerramento"
        p_novo = caixa_obrigado.text_frame.add_paragraph()
        p_novo.alignment = caixa_obrigado.text_frame.paragraphs[0].alignment
        run = p_novo.add_run()
        run.text = "{{NOME_CONCURSO}}"
        run.font.name = FONTE_UI_PPTX
        run.font.size = PPt(20)

    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    caminho = os.path.join(TEMPLATES_DIR, "slides.pptx")
    prs.save(caminho)
    return caminho


if __name__ == "__main__":
    print("DOCX (lista):", build_lista_template())
    print("DOCX (comentada):", build_comentada_template())
    print("DOCX (rastreabilidade):", build_rastreabilidade_template())
    print("PPTX:", build_pptx_template())

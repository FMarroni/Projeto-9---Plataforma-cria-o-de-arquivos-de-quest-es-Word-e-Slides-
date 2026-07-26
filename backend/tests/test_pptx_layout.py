import os
import re

from pptx import Presentation

from app.formatting import Token
from app.pptx_layout import (
    TAMANHO_FONTE_CORPO_PT,
    LayoutBlock,
    TextLayoutConfig,
    _simular_linhas,
    altura_linha_pt,
    largura_texto_pt,
    paginar_blocos,
)
from app.pptx_xml_utils import encontrar_shape_por_nome

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "templates", "slides.pptx")


def _config(largura: float = 670.0, altura: float = 320.0) -> TextLayoutConfig:
    return TextLayoutConfig(largura_util_pt=largura, altura_util_pt=altura, tamanho_fonte_pt=TAMANHO_FONTE_CORPO_PT)


def _altura_pagina(pagina_tokens: list[Token], config: TextLayoutConfig) -> float:
    """Reagrupa os tokens de uma página (já concatenados por `_anexar`, com
    `new_paragraph` marcando fronteira de bloco) em parágrafos e soma a altura
    simulada de cada um — usado para conferir se a página ficou dentro do
    orçamento vertical calculado."""
    paragrafos: list[list[Token]] = []
    atual: list[Token] = []
    for token in pagina_tokens:
        if token.new_paragraph:
            paragrafos.append(atual)
            atual = []
        else:
            atual.append(token)
    paragrafos.append(atual)
    return sum(len(_simular_linhas(p, config)) for p in paragrafos) * config.altura_linha_pt


# --- geometria real do template ----------------------------------------------


def test_text_layout_config_from_shape_usa_geometria_real_do_template():
    prs = Presentation(TEMPLATE_PATH)
    corpo = encontrar_shape_por_nome(prs.slides[2], "ph_corpo_questao")
    config = TextLayoutConfig.from_shape(corpo)

    # ~670pt de largura útil e ~300-320pt de altura útil (ver comentário
    # histórico em pptx_gen.py) — não são mais números repetidos no código,
    # vêm da geometria real do shape menos as margens internas do text_frame.
    assert 600 < config.largura_util_pt < 700
    assert 250 < config.altura_util_pt < 330
    assert config.altura_linha_pt > 0


# --- 1) medição tipográfica real, não contagem de caracteres -----------------


def test_largura_distingue_caracteres_estreitos_de_largos_com_mesma_contagem():
    estreito = largura_texto_pt("i" * 16, False, False, TAMANHO_FONTE_CORPO_PT)
    largo = largura_texto_pt("W" * 16, False, False, TAMANHO_FONTE_CORPO_PT)
    assert estreito > 0
    assert largo > estreito * 3  # mesmíssimo nº de caracteres, largura bem diferente


def test_largura_negrito_maior_que_regular():
    regular = largura_texto_pt("Texto de exemplo", False, False, TAMANHO_FONTE_CORPO_PT)
    negrito = largura_texto_pt("Texto de exemplo", True, False, TAMANHO_FONTE_CORPO_PT)
    assert negrito > regular


def test_largura_espaco_nao_e_zero_mas_vazio_e():
    assert largura_texto_pt(" ", False, False, TAMANHO_FONTE_CORPO_PT) > 0
    assert largura_texto_pt("", False, False, TAMANHO_FONTE_CORPO_PT) == 0.0


def test_altura_linha_escala_linearmente_com_o_tamanho_da_fonte():
    a16 = altura_linha_pt(16.0)
    a32 = altura_linha_pt(32.0)
    assert a16 > 0
    assert abs(a32 - a16 * 2) < 1e-6


# --- 2) vários parágrafos curtos aproveitam a área do slide ------------------


def test_varios_paragrafos_curtos_cabem_juntos_numa_pagina():
    config = _config()
    blocos = [
        LayoutBlock(tokens=[Token(text=f"Frase curta número {n}.")], keep_together=False) for n in range(6)
    ]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) == 1


# --- 3) parágrafo único maior que um slide é fragmentado ---------------------


def test_paragrafo_maior_que_um_slide_e_fragmentado_em_varias_paginas():
    config = _config(altura=100.0)
    texto_longo = " ".join(["palavra"] * 300)
    blocos = [LayoutBlock(tokens=[Token(text=texto_longo)], keep_together=False)]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) > 1
    reconstruido = " ".join("".join(t.text for t in p) for p in paginas)
    assert reconstruido.split().count("palavra") == 300


# --- 4) sequência extremamente longa sem espaços -----------------------------


def test_sequencia_sem_espacos_fatiada_por_caractere_sem_perder_texto():
    config = _config(largura=200.0)
    sequencia = "a" * 1000
    linhas = _simular_linhas([Token(text=sequencia)], config)
    assert len(linhas) > 1
    reconstruido = "".join(t.text for linha in linhas for t in linha)
    assert reconstruido == sequencia


# --- 5) preservação de estilos após fragmentação -----------------------------


def test_estilos_bold_italic_underline_preservados_ao_fatiar_por_caractere():
    config = _config(largura=150.0)
    token = Token(text="x" * 200, bold=True, italic=True, underline=True)
    fragmentos = [f for linha in _simular_linhas([token], config) for f in linha]
    assert len(fragmentos) > 1
    assert all(f.bold and f.italic and f.underline for f in fragmentos)
    assert "".join(f.text for f in fragmentos) == token.text


def test_subscrito_e_sobrescrito_preservados_no_word_wrap():
    config = _config()
    tokens = [
        Token(text="x"),
        Token(text="2", superscript=True),
        Token(text=" normal "),
        Token(text="i", subscript=True),
    ]
    fragmentos = [f for linha in _simular_linhas(tokens, config) for f in linha]
    superscritos = [f for f in fragmentos if f.superscript]
    subscritos = [f for f in fragmentos if f.subscript]
    assert superscritos and all(f.text == "2" for f in superscritos)
    assert subscritos and all(f.text == "i" for f in subscritos)


# --- 6/7) alternativas: inteiras quando cabem, fragmentadas quando maiores ---


def test_alternativa_curta_fica_inteira_numa_pagina_vazia():
    config = _config()
    blocos = [
        LayoutBlock(tokens=[Token(text="a) Alternativa curta.")], keep_together=True),
        LayoutBlock(tokens=[Token(text="b) Outra alternativa curta.")], keep_together=True),
    ]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) == 1


def test_alternativa_maior_que_slide_e_fragmentada_com_seguranca():
    config = _config(altura=100.0)
    texto_gigante = "a) " + " ".join(["palavra"] * 300)
    blocos = [LayoutBlock(tokens=[Token(text=texto_gigante)], keep_together=True)]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) > 1
    reconstruido = " ".join("".join(t.text for t in p) for p in paginas)
    assert reconstruido.split().count("palavra") == 300
    for pagina in paginas:
        assert _altura_pagina(pagina, config) <= config.altura_util_pt + 1e-6


# --- 8) imagem/equação sempre isolada em página própria ----------------------


def test_imagem_sempre_isolada_em_pagina_propria():
    config = _config()
    blocos = [
        LayoutBlock(tokens=[Token(text="antes")], keep_together=False),
        LayoutBlock(is_image=True, image_id="IMAGEM_01"),
        LayoutBlock(tokens=[Token(text="depois")], keep_together=False),
    ]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) == 3
    assert paginas[1] == [Token(text="", is_image=True, image_id="IMAGEM_01")]


# --- 9) parágrafos vazios continuam consumindo altura real -------------------


def test_paragrafo_vazio_conta_como_uma_linha_de_altura():
    config = _config()
    linhas = _simular_linhas([], config)
    assert linhas == [[]]


def test_bloco_vazio_entre_conteudos_desloca_o_que_vem_depois():
    linha = altura_linha_pt(TAMANHO_FONTE_CORPO_PT)
    config = _config(altura=linha * 2.5)  # cabem só 2 linhas cheias por página
    blocos = [
        LayoutBlock(tokens=[Token(text="Uma linha.")], keep_together=False),
        LayoutBlock(tokens=[]),  # linha em branco -- se não consumisse altura, tudo caberia numa só página
        LayoutBlock(tokens=[Token(text="Outra linha.")], keep_together=False),
    ]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) == 2


# --- 10) nunca gera página de conteúdo vazia ---------------------------------


def test_bloco_totalmente_vazio_nao_gera_pagina():
    config = _config()
    paginas = paginar_blocos([LayoutBlock(tokens=[])], config)
    assert paginas == []


def test_bloco_vazio_narrativo_no_limite_da_pagina_nunca_gera_pagina_em_branco():
    """Regressão: uma linha em branco dentro de um parágrafo narrativo (ex.:
    quebra dupla no meio do comentário, virando um LayoutBlock(tokens=[],
    keep_together=False)) podia cair exatamente no fim de uma página quase
    cheia e ser roteada para o caminho de fragmentação -- como um bloco vazio
    não tem fragmento nenhum para escrever, isso gerava uma página de
    conteúdo totalmente em branco entre dois parágrafos reais."""
    linha = altura_linha_pt(TAMANHO_FONTE_CORPO_PT)
    # largura estreita o bastante para cada palavra ocupar sua própria linha
    # (largura de "palavraN" ~69pt) e altura que deixa só ~0,5 linha livre
    # depois do primeiro bloco de 3 linhas -- exatamente o cenário do bug.
    config = TextLayoutConfig(largura_util_pt=75.0, altura_util_pt=linha * 3.5, tamanho_fonte_pt=TAMANHO_FONTE_CORPO_PT)
    blocos = [
        LayoutBlock(tokens=[Token(text="palavra0 palavra1 palavra2")], keep_together=False),
        LayoutBlock(tokens=[], keep_together=False),  # linha em branco dentro do comentário
        LayoutBlock(tokens=[Token(text="ultima")], keep_together=False),
    ]

    paginas = paginar_blocos(blocos, config)

    assert all(pagina for pagina in paginas), f"encontrada página vazia em: {paginas}"
    texto = "".join(t.text for p in paginas for t in p)
    for palavra in ("palavra0", "palavra1", "palavra2", "ultima"):
        assert palavra in texto


def test_paginar_blocos_nao_produz_nenhuma_pagina_vazia_em_caso_geral():
    config = _config(altura=90.0)
    blocos = [
        LayoutBlock(tokens=[Token(text=" ".join(["palavra"] * n))], keep_together=False) for n in [3, 50, 1, 70]
    ]
    paginas = paginar_blocos(blocos, config)
    assert all(pagina for pagina in paginas)


# --- 13) narrativa longa aproveita o espaço restante depois de rótulos ------


def test_narrativa_longa_ocupa_espaco_restante_apos_rotulo_curto():
    config = _config()
    blocos = [
        LayoutBlock(tokens=[Token(text="Gabarito: A", bold=True)], keep_together=True, keep_with_next=True),
        LayoutBlock(tokens=[]),
        LayoutBlock(tokens=[Token(text="Comentário:", bold=True)], keep_together=True, keep_with_next=True),
        LayoutBlock(tokens=[Token(text=" ".join(["explicação"] * 100))], keep_together=False),
    ]
    paginas = paginar_blocos(blocos, config)
    primeira_pagina_texto = "".join(t.text for t in paginas[0])
    assert "Gabarito: A" in primeira_pagina_texto
    assert "Comentário:" in primeira_pagina_texto
    # o rótulo não pode ficar sozinho: pelo menos o início da explicação
    # precisa estar na MESMA página que os rótulos
    assert "explicação" in primeira_pagina_texto


# --- 14) nenhuma perda, duplicação ou reordenação de texto -------------------


def test_paginacao_preserva_todo_o_texto_sem_perda_duplicacao_ou_reordenacao():
    config = _config(altura=150.0)
    paragrafos_originais = [
        "Primeiro parágrafo com texto normal para o teste de integridade.",
        "Segundo parágrafo, mais longo, com várias palavras diferentes para preencher espaço razoável.",
        "",  # parágrafo vazio no meio
        "a) Alternativa com texto suficiente para ocupar espaço razoável dentro do slide.",
        "b) " + " ".join(f"item{n}" for n in range(120)),  # força fragmentação
    ]
    blocos = [
        LayoutBlock(tokens=[Token(text=p)] if p else [], keep_together=(i >= 3))
        for i, p in enumerate(paragrafos_originais)
    ]

    paginas = paginar_blocos(blocos, config)
    texto_reconstruido = " ".join(t.text for pagina in paginas for t in pagina if not t.is_image)
    texto_normalizado = re.sub(r"\s+", " ", texto_reconstruido).strip()

    for p in paragrafos_originais:
        if p:
            assert re.sub(r"\s+", " ", p).strip() in texto_normalizado

    # nenhuma duplicação: cada "itemN" da alternativa gigante aparece 1x só
    for n in range(120):
        assert texto_reconstruido.count(f"item{n} ") + texto_reconstruido.count(f"item{n}\n") == 1 or (
            f"item{n}" in texto_reconstruido and texto_reconstruido.count(f"item{n}") == 1
        )


# --- 17) todas as páginas respeitam o orçamento vertical calculado ----------


def test_todas_as_paginas_respeitam_o_orcamento_vertical():
    config = _config(altura=200.0)
    blocos = [
        LayoutBlock(tokens=[Token(text=" ".join(["palavra"] * n))], keep_together=False)
        for n in [5, 40, 1, 60, 3, 80, 200]
    ]
    paginas = paginar_blocos(blocos, config)
    for pagina in paginas:
        if pagina and pagina[0].is_image:
            continue
        assert _altura_pagina(pagina, config) <= config.altura_util_pt + 1e-6


# --- fórmula com altura natural: flui com o texto, não força página própria -


def test_formula_com_altura_natural_flui_com_texto_antes_e_depois():
    config = _config()
    linha = config.altura_linha_pt
    blocos = [
        LayoutBlock(tokens=[Token(text="antes")], keep_together=False),
        LayoutBlock(is_image=True, image_id="EQUACAO_abc", altura_natural_pt=linha * 2),
        LayoutBlock(tokens=[Token(text="depois")], keep_together=False),
    ]
    paginas = paginar_blocos(blocos, config)
    # tudo cabe numa única página -- a fórmula não força mais um slide à parte
    assert len(paginas) == 1
    tokens_pagina = paginas[0]
    assert any(t.is_image and t.image_id == "EQUACAO_abc" for t in tokens_pagina)
    assert any(t.text == "antes" for t in tokens_pagina)
    assert any(t.text == "depois" for t in tokens_pagina)


def test_formula_com_altura_natural_avanca_pagina_quando_nao_cabe():
    linha_altura = 20.0
    config = TextLayoutConfig(largura_util_pt=670.0, altura_util_pt=linha_altura * 3, tamanho_fonte_pt=16.0)
    # preenche quase toda a página com texto antes da fórmula
    blocos = [
        LayoutBlock(tokens=[Token(text="x")], keep_together=False),
        LayoutBlock(is_image=True, image_id="EQUACAO_grande", altura_natural_pt=linha_altura * 2.9),
    ]
    # sem espaço nenhum sobrando após o texto consumir quase tudo — a fórmula
    # deve ser adiada para uma nova página, nunca fragmentada.
    paginas = paginar_blocos(blocos, config)
    pagina_com_formula = next(p for p in paginas if any(t.is_image for t in p))
    assert len(pagina_com_formula) == 1  # sozinha na sua página (não coube com "x")
    assert pagina_com_formula[0].image_id == "EQUACAO_grande"


def test_formula_de_conteudo_sem_altura_natural_continua_isolada():
    # regressão: imagem de CONTEÚDO comum (altura_natural_pt=None, o padrão)
    # preserva o comportamento histórico -- nunca divide slide com texto,
    # mesmo que o texto ao redor seja mínimo.
    config = _config()
    blocos = [
        LayoutBlock(tokens=[Token(text="antes")], keep_together=False),
        LayoutBlock(is_image=True, image_id="IMAGEM_01"),
        LayoutBlock(tokens=[Token(text="depois")], keep_together=False),
    ]
    paginas = paginar_blocos(blocos, config)
    assert len(paginas) == 3
    assert paginas[1] == [Token(text="", is_image=True, image_id="IMAGEM_01")]


def test_formula_maior_que_pagina_vazia_inteira_nao_e_fragmentada():
    config = _config(altura=50.0)
    blocos = [LayoutBlock(is_image=True, image_id="EQUACAO_enorme", altura_natural_pt=500.0)]
    paginas = paginar_blocos(blocos, config)
    # não fragmenta uma imagem -- insere inteira mesmo excedendo o orçamento
    assert len(paginas) == 1
    assert paginas[0] == [Token(text="", is_image=True, image_id="EQUACAO_enorme")]

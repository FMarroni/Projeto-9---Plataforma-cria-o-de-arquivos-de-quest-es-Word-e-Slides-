import io
import os
import re

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation

from app.formatting import Token
from app.formula_resolve import resolver_formulas_questao
from app.pptx_gen import (
    FONTE_CORPO,
    TAMANHO_FONTE_CORPO_PT,
    _agrupar_em_paragrafos,
    _agrupar_por_concurso_banca,
    _construir_blocos_enunciado,
    gerar_pptx,
)
from app.pptx_layout import LayoutBlock, TextLayoutConfig, paginar_blocos
from app.pptx_xml_utils import COR_MARCA, encontrar_shape_por_nome
from app.schemas import Alternativa, ExtractionResult, Formula, Questao

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "..", "templates", "slides.pptx")


def _questao(**overrides):
    base = dict(
        numero=1,
        banca="FCC",
        ano=2025,
        materia="Direito Penal",
        assunto="Furto",
        enunciado="Enunciado da questão.",
        alternativas=[Alternativa(letra="a", texto="Certo"), Alternativa(letra="b", texto="Errado")],
        gabarito="A",
        comentario="Comentário da questão.",
    )
    base.update(overrides)
    return Questao(**base)


def _texto_png() -> bytes:
    """PNG mínimo (não é um placeholder vazio) para testar o caminho de
    imagem sem depender de nenhum arquivo de amostra."""
    fig = plt.figure(figsize=(1, 1))
    plt.text(0.5, 0.5, "fig")
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    plt.close(fig)
    return buf.getvalue()


def _blocos_texto_por_slide(prs: Presentation) -> list[str]:
    """Texto do `ph_corpo_questao` de cada slide que tem esse shape (None nos
    demais — capa/divisor/encerramento não têm)."""
    textos = []
    for slide in prs.slides:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        textos.append(corpo.text_frame.text if corpo is not None else None)
    return textos


# --- _agrupar_em_paragrafos (inalterado — ainda usado por _blocos_de) --------


def test_agrupar_em_paragrafos_preserva_parágrafos_vazios():
    tokens = [
        Token(text="a"),
        Token(text="", new_paragraph=True),
        Token(text="", new_paragraph=True),  # parágrafo vazio no meio
        Token(text="b"),
    ]
    unidades = _agrupar_em_paragrafos(tokens)
    assert unidades == [[Token(text="a")], [], [Token(text="b")]]


# --- construção de blocos: sem gabarito/comentário; cor de marca no enunciado -


def _texto_dos_blocos(blocos: list[LayoutBlock]) -> str:
    return "".join(t.text for b in blocos for t in b.tokens)


def test_gabarito_e_comentario_nao_aparecem_nos_blocos_do_slide():
    # os slides são material de apoio para aula ao vivo -- gabarito/comentário
    # continuam nos DOCX, nunca nos blocos que viram slide de PPTX.
    questao = _questao()
    blocos = _construir_blocos_enunciado(questao)
    texto = _texto_dos_blocos(blocos)
    assert "Gabarito" not in texto
    assert "Comentário" not in texto


def test_alternativas_sao_atomicas_enunciado_e_narrativo():
    questao = _questao(enunciado="Enunciado longo o suficiente.")
    blocos = _construir_blocos_enunciado(questao)
    bloco_enunciado = blocos[0]
    blocos_alternativas = [b for b in blocos if b.tokens and re.match(r"^[ab]\)", b.tokens[0].text)]
    assert not bloco_enunciado.keep_together
    assert len(blocos_alternativas) == 2
    assert all(b.keep_together for b in blocos_alternativas)


def test_enunciado_marcado_com_cor_de_marca_e_negrito_alternativas_nao():
    questao = _questao(enunciado="Enunciado com *itálico* misturado.")
    blocos = _construir_blocos_enunciado(questao)

    bloco_enunciado = blocos[0]
    assert bloco_enunciado.tokens  # não vazio
    assert all(t.cor_marca and t.bold for t in bloco_enunciado.tokens)

    blocos_alternativas = [b for b in blocos if b.tokens and re.match(r"^[ab]\)", b.tokens[0].text)]
    assert blocos_alternativas
    for bloco in blocos_alternativas:
        assert not any(t.cor_marca for t in bloco.tokens)


# --- agrupamento por concurso/banca (banca + órgão + cargo + ano) -----------


def test_agrupar_por_concurso_banca_junta_mesma_chave_e_separa_diferentes():
    q1 = _questao(numero=1, banca="CEBRASPE", orgao="ANM", cargo="AA (ANM)", ano=2025)
    q2 = _questao(numero=2, banca="CEBRASPE", orgao="ANM", cargo="AA (ANM)", ano=2025)  # mesma chave de q1
    q3 = _questao(numero=3, banca="CEBRASPE", orgao="InoversaSul", cargo="Ana (InoversaSul)", ano=2025)  # órgão/cargo diferentes
    q4 = _questao(numero=4, banca="CEBRASPE", orgao="ANM", cargo="AA (ANM)", ano=2020)  # ano diferente

    grupos = _agrupar_por_concurso_banca([q1, q2, q3, q4])

    assert len(grupos) == 3
    assert [q.numero for q in grupos[0]] == [1, 2]
    assert [q.numero for q in grupos[1]] == [3]
    assert [q.numero for q in grupos[2]] == [4]


def test_agrupar_por_concurso_banca_preserva_ordem_de_primeira_aparicao():
    q1 = _questao(numero=1, orgao="A")
    q2 = _questao(numero=2, orgao="B")
    q3 = _questao(numero=3, orgao="A")  # mesma chave de q1 -- entra no grupo "A", não cria um 3º grupo

    grupos = _agrupar_por_concurso_banca([q1, q2, q3])

    assert len(grupos) == 2  # só 2 chaves distintas (A e B), não 3 questões = 3 grupos
    assert [q.numero for q in grupos[0]] == [1, 3]  # grupo "A": ordem de aparição original preservada
    assert [q.numero for q in grupos[1]] == [2]  # grupo "B"


# --- geração completa via gerar_pptx (integração) ---------------------------


def _extraction_simples():
    return ExtractionResult(
        bancas=["FCC"],
        anos=[2025],
        cargos=["Analista"],
        questoes=[_questao(numero=1)],
    )


def test_ordem_capa_contracapa_divisor_concurso_banca_questao_encerramento(tmp_path):
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    layouts = [slide.slide_layout.name for slide in prs.slides]
    assert layouts[0] == "Capa"
    assert layouts[1] == "Contracapa"
    assert layouts[2] == "Capa_Disciplina"
    assert layouts[3] == "Capa_Concurso_Banca"
    assert layouts[4] == "Questao_Automatica"
    assert layouts[-1] == "Default 2"  # Encerramento, preservado do template


def test_capa_e_contracapa_nao_recebem_nenhuma_edicao(tmp_path):
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    for slide in (prs.slides[0], prs.slides[1]):
        assert not any(sh.has_text_frame and sh.text_frame.text.strip() for sh in slide.shapes)


def test_divisor_disciplina_mostra_nome_e_prof(tmp_path):
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    slide_divisor = prs.slides[2]
    ph_disciplina = encontrar_shape_por_nome(slide_divisor, "ph_disciplina")
    ph_prof = encontrar_shape_por_nome(slide_divisor, "ph_prof_disciplina")
    assert ph_disciplina is not None and "DIREITO PENAL" in ph_disciplina.text_frame.text
    assert ph_prof is not None and ph_prof.text_frame.text == "Prof."


def test_arquivo_gerado_nao_tem_partes_xml_duplicadas(tmp_path):
    # regressão: remover os slides-stencil obsoletos ANTES de inserir os
    # novos fazia o add_slide() reaproveitar um nome de part (slideN.xml)
    # ainda ocupado pelo slide de Encerramento, corrompendo o arquivo salvo.
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    nomes = [part.partname for part in prs.part.package.iter_parts() if "/ppt/slides/slide" in part.partname]
    assert len(nomes) == len(set(nomes)), f"partnames duplicados: {nomes}"


def test_slides_nao_trazem_gabarito_nem_comentario(tmp_path):
    # os slides são material de apoio para aula ao vivo -- gabarito/comentário
    # nunca aparecem neles, mesmo a questão tendo os dois preenchidos.
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    textos_por_slide = _blocos_texto_por_slide(prs)
    slides_com_enunciado = [t for t in textos_por_slide if t and "Enunciado da questão" in t]

    assert slides_com_enunciado, "esperava ao menos um slide de enunciado"
    for t in textos_por_slide:
        if t:
            assert "Gabarito" not in t
            assert "Comentário" not in t


def test_questao_longa_gera_mais_de_um_slide_de_enunciado_sem_estourar_orcamento(tmp_path):
    enunciado_longo = "\n".join(
        f"Frase número {n} bem detalhada para ocupar bastante espaço na caixa de corpo." for n in range(40)
    )
    questao = _questao(numero=1, enunciado=enunciado_longo, alternativas=[])
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])

    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    slides_com_enunciado = [
        slide
        for slide in prs.slides
        if (corpo := encontrar_shape_por_nome(slide, "ph_corpo_questao")) is not None
        and "Frase número" in (corpo.text_frame.text or "")
    ]

    assert len(slides_com_enunciado) > 1


# --- critérios de aceitação: sem autofit, fonte fixa, sem transbordo --------


def test_nenhum_slide_usa_text_to_fit_shape_ou_normautofit(tmp_path):
    saida = gerar_pptx(_extraction_simples(), TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    for slide in prs.slides:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        if corpo is None:
            continue
        assert corpo.text_frame.auto_size is None
        assert corpo.text_frame._txBody.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
        ).find("{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit") is None


def test_todos_os_runs_do_corpo_tem_fonte_e_tamanho_explicitos(tmp_path):
    questao = _questao(
        enunciado="Enunciado com *itálico*, **negrito** e $x^2$ para variar o estilo dos runs.",
    )
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    algum_run_verificado = False
    for slide in prs.slides:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        if corpo is None:
            continue
        for paragrafo in corpo.text_frame.paragraphs:
            for run in paragrafo.runs:
                assert run.font.name == FONTE_CORPO
                assert run.font.size is not None
                assert round(run.font.size.pt, 3) == TAMANHO_FONTE_CORPO_PT
                algum_run_verificado = True

    assert algum_run_verificado


def test_nenhum_slide_de_conteudo_fica_vazio(tmp_path):
    enunciado_longo = " ".join(f"palavra{n}" for n in range(300))
    questao = _questao(enunciado=enunciado_longo, alternativas=[])
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    for slide in prs.slides:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        if corpo is None:
            continue  # capa/divisor/encerramento não têm esse shape
        tem_texto = bool((corpo.text_frame.text or "").strip())
        tem_imagem = any(sh.shape_type == 13 for sh in slide.shapes)  # PICTURE
        assert tem_texto or tem_imagem


def test_alternativa_curta_nunca_e_fragmentada(tmp_path):
    questao = _questao(
        enunciado="Enunciado curto.",
        alternativas=[Alternativa(letra=l, texto=f"Alternativa {l} bem curta.") for l in "abcde"],
    )
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    textos = _blocos_texto_por_slide(prs)
    for letra in "abcde":
        alternativa_completa = f"{letra}) Alternativa {letra} bem curta."
        assert any(t and alternativa_completa in t for t in textos)


def test_alternativa_maior_que_slide_e_fragmentada_sem_perder_texto(tmp_path):
    texto_gigante = " ".join(f"palavra{n}" for n in range(400))
    questao = _questao(
        enunciado="Enunciado curto.",
        alternativas=[Alternativa(letra="a", texto=texto_gigante)],
    )
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    textos = " ".join(t for t in _blocos_texto_por_slide(prs) if t)
    # \b (fronteira de palavra): sem isso, "palavra7" também "casaria" como
    # substring dentro de "palavra70".."palavra79"
    for n in range(400):
        assert len(re.findall(rf"\bpalavra{n}\b", textos)) == 1


def test_imagem_de_conteudo_fica_em_pagina_propria(tmp_path):
    # imagem de CONTEÚDO (gráfico/tabela real do PDF, [IMAGEM_NN]) preserva o
    # comportamento histórico: grande o bastante para justificar uma página
    # exclusiva, nunca divide slide com texto.
    imagem = _texto_png()
    enunciado = "Considere [IMAGEM_01] o gráfico acima para responder."
    questao = _questao(enunciado=enunciado, alternativas=[])
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"), {"IMAGEM_01": imagem})
    prs = Presentation(saida)

    # restrito a slides de QUESTÃO (têm ph_header_questao) — capa/encerramento
    # do próprio template corporativo já trazem imagens decorativas, que não
    # têm nada a ver com o [IMAGEM_NN] desta questão.
    slides_de_questao_com_imagem = [
        slide
        for slide in prs.slides
        if encontrar_shape_por_nome(slide, "ph_header_questao") is not None
        and any(sh.shape_type == 13 for sh in slide.shapes)
    ]
    assert len(slides_de_questao_com_imagem) >= 1
    for slide in slides_de_questao_com_imagem:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        assert corpo is not None
        assert not (corpo.text_frame.text or "").strip()


def test_equacao_pequena_flui_com_texto_ao_redor_sem_orfanizar_slide(tmp_path):
    # requisito: uma fórmula pequena (EQUACAO_*/FORMULA_*) NÃO deve mais
    # forçar um slide "só fórmula" órfão quando o texto ao redor caberia
    # junto — diferente da imagem de conteúdo acima, que continua sozinha.
    enunciado = (
        "Considere a fórmula $$x^2 + y^2 = z^2$$ para responder à questão a seguir."
    )
    questao = _questao(enunciado=enunciado, alternativas=[])
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    slides_de_questao = [s for s in prs.slides if encontrar_shape_por_nome(s, "ph_header_questao") is not None]
    # a fórmula não pode ter gerado um slide adicional sozinho — a questão
    # inteira (texto curto + 1 fórmula pequena) cabe num único slide.
    assert len(slides_de_questao) == 1

    slide = slides_de_questao[0]
    tem_imagem = any(sh.shape_type == 13 for sh in slide.shapes)
    assert tem_imagem  # a equação foi de fato renderizada como imagem

    texto_do_slide = " ".join(
        (sh.text_frame.text or "") for sh in slide.shapes if sh.has_text_frame
    )
    assert "Considere a fórmula" in texto_do_slide
    assert "para responder à questão a seguir" in texto_do_slide


def test_geracao_nao_perde_duplica_ou_reordena_texto(tmp_path):
    enunciado_longo = " ".join(f"enun{n}" for n in range(150))
    questao = _questao(
        enunciado=enunciado_longo,
        alternativas=[Alternativa(letra=l, texto=f"alt{l}_{n}") for n, l in enumerate("abcde")],
    )
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    texto_total = " ".join(t for t in _blocos_texto_por_slide(prs) if t)

    # \b garante fronteira de palavra -- sem isso, "enun1" também "casaria"
    # como substring dentro de "enun10".."enun19"/"enun100".."enun149"
    for n in range(150):
        assert len(re.findall(rf"\benun{n}\b", texto_total)) == 1
    for l in "abcde":
        assert f"{l}) alt{l}_" in texto_total


# --- critério de aceitação: cor de marca no enunciado, cor padrão nas alternativas ---


def test_slide_gerado_aplica_cor_de_marca_so_no_enunciado(tmp_path):
    questao = _questao(
        enunciado="Enunciado que deve sair em roxo.",
        alternativas=[Alternativa(letra="a", texto="Alternativa em cor padrão.")],
    )
    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    runs_enunciado_roxos = []
    runs_alternativa_com_cor = []
    for slide in prs.slides:
        corpo = encontrar_shape_por_nome(slide, "ph_corpo_questao")
        if corpo is None:
            continue
        for paragrafo in corpo.text_frame.paragraphs:
            for run in paragrafo.runs:
                if "roxo" in run.text:
                    runs_enunciado_roxos.append(run)
                if "padrão" in run.text and run.font.color.type is not None:
                    runs_alternativa_com_cor.append(run)

    assert runs_enunciado_roxos, "esperava encontrar o run do enunciado"
    for run in runs_enunciado_roxos:
        assert run.font.color.rgb == COR_MARCA
        assert run.font.bold

    assert not runs_alternativa_com_cor, "alternativa não deveria ter cor explícita"


# --- hierarquia completa: disciplina > concurso/banca > questões ------------


def test_hierarquia_completa_disciplina_e_concurso_banca(tmp_path):
    questoes = [
        _questao(
            numero=1, banca="CEBRASPE", orgao="ANM", cargo="AA (ANM)", ano=2025,
            materia="Administração", enunciado="Questão um do primeiro subgrupo.",
        ),
        _questao(
            numero=2, banca="CEBRASPE", orgao="ANM", cargo="AA (ANM)", ano=2025,
            materia="Administração", enunciado="Questão dois, mesmo subgrupo da anterior.",
        ),
        _questao(
            numero=3, banca="CEBRASPE", orgao="InoversaSul", cargo="Ana (InoversaSul)", ano=2025,
            materia="Administração", enunciado="Questão três, novo subgrupo (órgão/cargo diferentes).",
        ),
        _questao(
            numero=4, banca="FCC", orgao="TRT4", cargo="AJ TRT4", ano=2022,
            materia="Matemática", enunciado="Questão quatro, nova disciplina.",
        ),
    ]
    extraction = ExtractionResult(
        bancas=["CEBRASPE", "FCC"], anos=[2022, 2025], cargos=["AA (ANM)", "Ana (InoversaSul)", "AJ TRT4"],
        questoes=questoes,
    )
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"))
    prs = Presentation(saida)

    layouts = [slide.slide_layout.name for slide in prs.slides]
    # Capa, Contracapa, Divisor(Administração), ConcursoBanca(ANM), Q1, Q2,
    # ConcursoBanca(InoversaSul), Q3, Divisor(Matemática), ConcursoBanca(TRT4), Q4, Encerramento
    assert layouts == [
        "Capa", "Contracapa",
        "Capa_Disciplina", "Capa_Concurso_Banca", "Questao_Automatica", "Questao_Automatica",
        "Capa_Concurso_Banca", "Questao_Automatica",
        "Capa_Disciplina", "Capa_Concurso_Banca", "Questao_Automatica",
        "Default 2",
    ]

    disciplinas = [
        encontrar_shape_por_nome(s, "ph_disciplina").text_frame.text
        for s in prs.slides if s.slide_layout.name == "Capa_Disciplina"
    ]
    assert disciplinas == ["ADMINISTRAÇÃO", "MATEMÁTICA"]

    concurso_banca_textos = [
        encontrar_shape_por_nome(s, "ph_concurso_banca").text_frame.text
        for s in prs.slides if s.slide_layout.name == "Capa_Concurso_Banca"
    ]
    assert concurso_banca_textos[0] == "CEBRASPE\nAA (ANM) — 2025"
    assert concurso_banca_textos[1] == "CEBRASPE\nAna (InoversaSul) — 2025"
    assert concurso_banca_textos[2] == "FCC\nAJ TRT4 — 2022"


# --- regressão ponta a ponta do bug relatado: G(t) = t³ - 23/2·t² + ... -----
#
# LaTeX exato do PDF original e a transcrição ERRADA já observada em produção
# (perde o termo cúbico "t^3" inteiro e reanexa o expoente "3" ao termo
# errado, "55/4 t^3" em vez de "55/4 t") — ver também test_formula_resolve.py.

LATEX_GT_CORRETO = r"G(t)=t^{3}-\frac{23}{2}t^{2}+\frac{55}{4}t+\frac{399}{8},\ t\in[0,10]"
LATEX_GT_ERRADO = r"G(t)=-\frac{23}{2}t^{2}+\frac{55}{4}t^{3}+\frac{399}{8}"


def test_gt_formula_correta_com_alta_confianca_gera_pptx_integro_com_imagem(tmp_path):
    enunciado = "Considere a função G definida por [FORMULA_01] no intervalo dado. Determine seu valor máximo."
    questao = _questao(
        enunciado=enunciado,
        alternativas=[],
        formulas=[Formula(id="FORMULA_01", latex=LATEX_GT_CORRETO, display=True, confidence=0.95)],
    )
    questao = resolver_formulas_questao(questao)
    assert "[FORMULA_01]" not in questao.enunciado  # aceita e substituída pelo LaTeX

    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"), {"FORMULA_01": _texto_png()})

    prs = Presentation(saida)  # não pode lançar -- arquivo íntegro
    slides_de_questao = [s for s in prs.slides if encontrar_shape_por_nome(s, "ph_header_questao") is not None]
    assert slides_de_questao
    assert any(any(sh.shape_type == 13 for sh in s.shapes) for s in slides_de_questao)  # a fórmula virou imagem

    _assert_pptx_nunca_contem_formula_corrompida(prs)


def test_gt_formula_com_baixa_confianca_usa_recorte_original_sem_corromper_pptx(tmp_path):
    # mesmo que a IA cometesse o erro real (perder o termo cúbico) mas
    # sinalizasse baixa confiança (comportamento honesto esperado), o
    # resultado final ainda é um PPTX íntegro com uma imagem (o recorte
    # original fiel) no lugar -- nunca o texto corrompido.
    enunciado = "Considere a função G definida por [FORMULA_01] no intervalo dado."
    questao = _questao(
        enunciado=enunciado,
        alternativas=[],
        formulas=[Formula(id="FORMULA_01", latex=LATEX_GT_ERRADO, display=True, confidence=0.2)],
    )
    questao = resolver_formulas_questao(questao)
    assert "[FORMULA_01]" in questao.enunciado  # rejeitada -- marcador intocado

    extraction = ExtractionResult(bancas=["FCC"], anos=[2025], cargos=["Analista"], questoes=[questao])
    saida = gerar_pptx(extraction, TEMPLATE_PATH, str(tmp_path / "slides.pptx"), {"FORMULA_01": _texto_png()})

    prs = Presentation(saida)
    slides_de_questao = [s for s in prs.slides if encontrar_shape_por_nome(s, "ph_header_questao") is not None]
    assert any(any(sh.shape_type == 13 for sh in s.shapes) for s in slides_de_questao)

    _assert_pptx_nunca_contem_formula_corrompida(prs)


def _assert_pptx_nunca_contem_formula_corrompida(prs: Presentation) -> None:
    """A fórmula errada (termo cúbico perdido/realocado) NUNCA pode aparecer
    como TEXTO literal em nenhum shape do arquivo final -- nos dois casos
    (LaTeX aceito ou recorte de fallback), ela sempre vira uma IMAGEM, nunca
    texto potencialmente corrompido gravado diretamente no XML do slide."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto = shape.text_frame.text
                assert LATEX_GT_ERRADO not in texto
                assert r"\frac{55}{4}t^{3}" not in texto

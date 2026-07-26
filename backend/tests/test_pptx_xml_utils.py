from pptx import Presentation
from pptx.util import Cm

from app.pptx_xml_utils import duplicar_linha_tabela, duplicar_slide, remover_slide


def _presentation_com_slide_de_texto():
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(2))
    box.name = "caixa_teste"
    box.text_frame.text = "conteúdo original"
    return prs


def test_duplicar_slide_aumenta_contagem_e_preserva_texto():
    prs = _presentation_com_slide_de_texto()
    assert len(prs.slides) == 1

    novo_idx = duplicar_slide(prs, 0)

    assert novo_idx == 1
    assert len(prs.slides) == 2
    texto_original = prs.slides[0].shapes[0].text_frame.text
    texto_copia = prs.slides[1].shapes[0].text_frame.text
    assert texto_original == texto_copia == "conteúdo original"

    # são shapes independentes (deepcopy), não a mesma referência de elemento XML
    assert prs.slides[0].shapes[0]._element is not prs.slides[1].shapes[0]._element


def test_duplicar_slide_com_insert_at_explicito_preserva_ordem():
    prs = _presentation_com_slide_de_texto()
    # duplica o slide 0 três vezes, sempre inserindo no fim (posição crescente),
    # simulando a paginação de várias páginas de uma mesma disciplina
    idx1 = duplicar_slide(prs, 0, insert_at=1)
    idx2 = duplicar_slide(prs, 0, insert_at=2)
    idx3 = duplicar_slide(prs, 0, insert_at=3)

    assert [idx1, idx2, idx3] == [1, 2, 3]
    assert len(prs.slides) == 4


def test_remover_slide():
    prs = _presentation_com_slide_de_texto()
    duplicar_slide(prs, 0)
    assert len(prs.slides) == 2

    remover_slide(prs, 0)
    assert len(prs.slides) == 1


def test_duplicar_linha_tabela_clona_e_permite_novo_conteudo():
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    shape = slide.shapes.add_table(2, 2, Cm(1), Cm(1), Cm(10), Cm(4))
    table = shape.table
    table.cell(0, 0).text = "cabecalho"
    table.cell(1, 0).text = "linha exemplo"

    assert len(table.rows) == 2

    nova_linha = duplicar_linha_tabela(table)
    assert len(table.rows) == 3

    nova_linha.cells[0].text = "linha nova"
    assert table.cell(2, 0).text == "linha nova"
    # a linha original de exemplo continua intacta
    assert table.cell(1, 0).text == "linha exemplo"

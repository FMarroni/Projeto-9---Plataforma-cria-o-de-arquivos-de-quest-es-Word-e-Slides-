"""Preenche `templates/slides.pptx` a partir de um ExtractionResult +
DisciplinaStats. Estrutura da apresentação (nomes dos LAYOUTS do slide
mestre — ver `templates/slides.pptx`, Modo de Exibição Mestre):

    1. "Capa"                — estático, sem edição (design já pronto).
    2. "Contracapa"          — estático, sem edição (design já pronto).
    3. Para cada disciplina:
       3.1. "Capa_Disciplina"      — divisor, mostra o nome da disciplina.
       3.2. Para cada subgrupo de concurso/banca (banca+órgão+cargo+ano):
            3.2.1. "Capa_Concurso_Banca" — divisor, mostra concurso/banca.
            3.2.2. "Questao_Automatica" (1+ por questão, ver paginação)
    4. Encerramento — ainda o slide-stencil pré-existente no template
       (`ph_encerramento`), preservado e só recebe substituição de token
       (nunca precisou de placeholder dinâmico por questão).

Diferente da versão anterior, os slides de Capa/Contracapa/Divisor de
disciplina/Divisor de concurso-banca NÃO são mais clonados de um
slide-stencil pré-existente — são criados direto a partir do LAYOUT do
slide mestre (`pptx_xml_utils.criar_slide_de_layout`), porque esses layouts
(autorados diretamente no PowerPoint, não mais derivados de um Google
Slides exportado) materializam os placeholders corretamente nesse caminho.
Os 3 slides-stencil obsoletos que ainda existirem no arquivo de template
(deixados por uma geração anterior a essa reestruturação) são descartados
no início — só o de Encerramento é preservado (ver `_preparar_template`).

Os slides de questão trazem SOMENTE enunciado + alternativas — sem gabarito
nem comentário (Épico "Sem Comentários/Gabarito nos slides"): o PPTX é
material de apoio para o professor dar aula ao vivo, não uma versão
"respondida" da prova; gabarito/comentário continuam nos DOCX (lista e
comentada), que são os documentos de estudo/revisão.

No corpo do slide, o enunciado sai na cor de marca (roxo, `Token.cor_marca`
— ver `pptx_xml_utils.COR_MARCA`) e em negrito, e as alternativas na cor
padrão — mesmo padrão visual usado nos DOCX (`docx_lista.py`/`docx_comentada.py`
via `cor=COR_MARCA, negrito_forcado=True` no cabeçalho do enunciado).

Paginação (ver `app.pptx_layout` para a implementação): cada documento vira
uma lista de `LayoutBlock`s semânticos (um por parágrafo/alternativa/rótulo,
mais um por imagem/equação já resolvida) e é distribuída entre slides por
`pptx_layout.paginar_blocos`, que mede a largura/altura REAL de cada trecho de
texto (métrica de glifo da própria fonte Montserrat gravada no PPTX) em vez de
contar caracteres — por isso um parágrafo maior que um slide inteiro é
fragmentado (nunca inserido sabendo que vai transbordar), e o corpo NÃO
depende mais de "encolher texto até caber" (`MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE`)
como rede de segurança.
"""

import os
from dataclasses import replace

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Emu, Pt

from app.analysis import DisciplinaStats, build_disciplinas_stats
from app.formatting import Token, tokenize_rich_text
from app.latex_render import resolver_tokens_equacao
from app.pptx_layout import (
    TAMANHO_FONTE_CORPO_PT,
    FONTE_CORPO,
    MARGEM_SEGURANCA_VERTICAL,
    LayoutBlock,
    TextLayoutConfig,
    _simular_linhas,
    paginar_blocos,
)
from app.pptx_xml_utils import (
    adicionar_imagem_slide,
    criar_slide_de_layout,
    dimensoes_nativas_imagem,
    dimensoes_proporcionais,
    encontrar_layout_por_nome,
    encontrar_shape_por_nome,
    escrever_texto_formatado,
    remover_slide,
    substituir_texto_mantendo_formatacao,
)
from app.schemas import ExtractionResult, Questao

# Nomes dos layouts no slide mestre (ver docstring do módulo).
LAYOUT_CAPA = "Capa"
LAYOUT_CONTRACAPA = "Contracapa"
LAYOUT_DIVISOR_DISCIPLINA = "Capa_Disciplina"
LAYOUT_DIVISOR_CONCURSO_BANCA = "Capa_Concurso_Banca"
LAYOUT_QUESTAO = "Questao_Automatica"

IDX_HEADER_QUESTAO = 10
IDX_CORPO_QUESTAO = 11
IDX_DISCIPLINA_TITULO = 13
IDX_DISCIPLINA_PROF = 14

# Geometria (EMU) e estilo do texto de concurso/banca em "Capa_Concurso_Banca"
# — esse layout ainda não tem um placeholder de texto de verdade ali (só uma
# caixa estática no design), então a caixa é recriada por código na mesma
# posição/estilo lidos diretamente do layout (ver conversa/investigação —
# `Google Shape;75;p18`), em vez de preencher um placeholder inexistente.
CONCURSO_BANCA_LEFT = Emu(949375)
CONCURSO_BANCA_TOP = Emu(1610675)
CONCURSO_BANCA_WIDTH = Emu(7684200)
CONCURSO_BANCA_HEIGHT = Emu(658601)
COR_CONCURSO_BANCA = RGBColor(0x42, 0x2E, 0xA4)  # cor exata lida do design (endParaRPr da caixa)
TAMANHO_CONCURSO_BANCA = Pt(32)

# Estilo (lido do design) do placeholder "Prof." em Capa_Disciplina — run-level
# explícito no layout original, por isso replicado explicitamente aqui em vez
# de confiar em herança de estilo (ver conversa/investigação).
TAMANHO_PROF_DISCIPLINA = Pt(16)

TAMANHO_FONTE_CORPO = Pt(TAMANHO_FONTE_CORPO_PT)

EQUACAO_MAX_LARGURA = Cm(14)
EQUACAO_MAX_ALTURA = Cm(6)


def _substituir_tokens_no_slide(slide, tokens: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for token, valor in tokens.items():
            substituir_texto_mantendo_formatacao(shape.text_frame, token, valor or "")


def _localizar_slides_pre_existentes(prs: Presentation) -> tuple[int, list[int]]:
    """Localiza, no template CARREGADO (antes de qualquer inserção), o índice
    do slide de Encerramento e os índices dos slides-stencil obsoletos
    deixados por uma geração anterior a esta reestruturação (Capa/Divisor/
    Questão antigos — agora nascem direto dos layouts, ver
    `criar_slide_de_layout`) — identificados pelo shape `ph_encerramento`
    (ou a ausência dele), não por índice fixo, para não depender da ordem
    exata deixada no arquivo.

    Importante: os stencils obsoletos só podem ser removidos DEPOIS de todas
    as inserções (ver `gerar_pptx`), nunca antes — removê-los primeiro faz o
    `add_slide()` do python-pptx reaproveitar o nome de part XML (`slideN.xml`)
    que o slide de Encerramento ainda ocupa, corrompendo o arquivo (duas
    partes com o mesmo nome dentro do .pptx)."""
    indice_encerramento = next(
        i for i, slide in enumerate(prs.slides) if encontrar_shape_por_nome(slide, "ph_encerramento") is not None
    )
    indices_obsoletos = [i for i in range(len(prs.slides)) if i != indice_encerramento]
    return indice_encerramento, indices_obsoletos


def _inserir_slide_estatico(prs: Presentation, layout, cursor: int) -> int:
    """Insere um slide sem nenhuma edição (Capa/Contracapa) — o design já
    está pronto no próprio layout do slide mestre."""
    criar_slide_de_layout(prs, layout, cursor)
    return cursor + 1


def _inserir_divisor_disciplina(prs: Presentation, layout, disciplina: DisciplinaStats, cursor: int) -> int:
    """Insere o slide-divisor com o nome da disciplina."""
    slide = criar_slide_de_layout(
        prs,
        layout,
        cursor,
        renomear_placeholders={IDX_DISCIPLINA_TITULO: "ph_disciplina", IDX_DISCIPLINA_PROF: "ph_prof_disciplina"},
    )

    ph_disciplina = encontrar_shape_por_nome(slide, "ph_disciplina")
    escrever_texto_formatado(ph_disciplina.text_frame, tokenize_rich_text(disciplina.nome.upper()))

    ph_prof = encontrar_shape_por_nome(slide, "ph_prof_disciplina")
    escrever_texto_formatado(ph_prof.text_frame, [Token(text="Prof.")])
    for paragrafo in ph_prof.text_frame.paragraphs:
        for run in paragrafo.runs:
            run.font.name = FONTE_CORPO
            run.font.size = TAMANHO_PROF_DISCIPLINA

    return cursor + 1


def _chave_concurso_banca(questao: Questao) -> tuple[str, str | None, str | None, int | None]:
    """Identifica um subgrupo de 'concurso/banca' por banca + órgão + cargo +
    ano — a mesma banca/órgão com cargo ou ano diferentes conta como um
    concurso distinto (critério confirmado com o usuário)."""
    return (questao.banca, questao.orgao, questao.cargo, questao.ano)


def _agrupar_por_concurso_banca(questoes: list[Questao]) -> list[list[Questao]]:
    """Agrupa `questoes` (já na ordem de uma disciplina) por
    `_chave_concurso_banca`, preservando a ordem de primeira aparição de cada
    grupo — nunca reordena as questões em si."""
    grupos: dict[tuple, list[Questao]] = {}
    ordem: list[tuple] = []
    for questao in questoes:
        chave = _chave_concurso_banca(questao)
        if chave not in grupos:
            grupos[chave] = []
            ordem.append(chave)
        grupos[chave].append(questao)
    return [grupos[chave] for chave in ordem]


def _texto_divisor_concurso_banca(questao: Questao) -> tuple[str, str]:
    """(linha 1, linha 2) exibidas no divisor de concurso/banca — banca em
    cima, cargo+ano (o que efetivamente diferencia o concurso) embaixo."""
    linha1 = questao.banca or ""
    partes = [p for p in (questao.cargo or questao.orgao, str(questao.ano) if questao.ano else None) if p]
    linha2 = " — ".join(partes)
    return linha1, linha2


def _inserir_divisor_concurso_banca(prs: Presentation, layout, questao_exemplo: Questao, cursor: int) -> int:
    """Insere o slide-divisor do subgrupo de concurso/banca — o layout ainda
    não tem um placeholder de texto pronto para isso (só uma caixa estática
    no design), então a caixa é recriada por código na posição/estilo lidos
    do layout (ver constantes `CONCURSO_BANCA_*`/`COR_CONCURSO_BANCA`)."""
    slide = criar_slide_de_layout(prs, layout, cursor)

    caixa = slide.shapes.add_textbox(
        CONCURSO_BANCA_LEFT, CONCURSO_BANCA_TOP, CONCURSO_BANCA_WIDTH, CONCURSO_BANCA_HEIGHT
    )
    caixa.name = "ph_concurso_banca"
    caixa.text_frame.word_wrap = True

    linha1, linha2 = _texto_divisor_concurso_banca(questao_exemplo)
    paragrafo1 = caixa.text_frame.paragraphs[0]
    paragrafo1.text = linha1
    paragrafo2 = caixa.text_frame.add_paragraph() if linha2 else None
    if paragrafo2 is not None:
        paragrafo2.text = linha2

    for paragrafo in caixa.text_frame.paragraphs:
        for run in paragrafo.runs:
            run.font.name = FONTE_CORPO
            run.font.size = TAMANHO_CONCURSO_BANCA
            run.font.bold = True
            run.font.color.rgb = COR_CONCURSO_BANCA

    return cursor + 1


def _agrupar_em_paragrafos(tokens: list[Token]) -> list[list[Token]]:
    """Agrupa a lista de Tokens em unidades de parágrafo: cada unidade vai até
    (sem incluir) o próximo Token `new_paragraph=True`, que passa a ser só o
    delimitador entre unidades — inclusive parágrafos vazios (linhas de
    espaçamento) viram uma unidade `[]` própria, contada como 1 linha real na
    paginação (ver `pptx_layout._simular_linhas`) em vez de ignorada. Um Token
    de imagem sempre fica isolado na própria unidade."""
    unidades: list[list[Token]] = []
    atual: list[Token] = []
    for token in tokens:
        if token.is_image:
            unidades.append(atual)
            atual = []
            unidades.append([token])
            continue
        if token.new_paragraph:
            unidades.append(atual)
            atual = []
            continue
        atual.append(token)
    unidades.append(atual)
    return unidades


def _blocos_de(tokens: list[Token], *, keep_together: bool) -> list[LayoutBlock]:
    """Agrupa `tokens` em parágrafos (`_agrupar_em_paragrafos`) e embrulha cada
    um num `LayoutBlock` com a flag `keep_together` dada — imagens (já
    isoladas por `_agrupar_em_paragrafos`) viram blocos `is_image=True`."""
    blocos: list[LayoutBlock] = []
    for unidade in _agrupar_em_paragrafos(tokens):
        if unidade and unidade[0].is_image:
            blocos.append(LayoutBlock(is_image=True, image_id=unidade[0].image_id))
        else:
            blocos.append(LayoutBlock(tokens=unidade, keep_together=keep_together))
    return blocos


def _tokens_enunciado(texto: str) -> list[Token]:
    """Tokeniza o enunciado marcando cada token com `cor_marca=True` (roxo) e
    `bold=True` — mesmo padrão visual do cabeçalho do enunciado nos DOCX
    (`docx_lista.py`/`docx_comentada.py`, via `cor=COR_MARCA,
    negrito_forcado=True`), só que aplicado por token (não por chamada,
    ver `Token.cor_marca`) porque um slide de pptx é escrito de uma vez só,
    podendo misturar a cauda do enunciado com o início das alternativas."""
    return [replace(t, cor_marca=True, bold=True) for t in tokenize_rich_text(texto)]


def _construir_blocos_enunciado(questao: Questao) -> list[LayoutBlock]:
    """Enunciado + alternativas — sem gabarito nem comentário (os slides são
    material de apoio para aula ao vivo, não uma versão "respondida" da
    prova; gabarito/comentário continuam nos DOCX). O corpo do enunciado é
    narrativo (`keep_together=False`): pode ser fatiado para aproveitar o
    espaço do slide, e sai em roxo/negrito (`_tokens_enunciado`). Cada
    alternativa é atômica (`keep_together=True`) e na cor padrão (sem
    `cor_marca`) — só é fragmentada se for maior que um slide inteiro; do
    contrário, uma alternativa inteira é preservada sempre que couber num
    slide vazio."""
    blocos = _blocos_de(_tokens_enunciado(questao.enunciado.strip()), keep_together=False)

    if questao.alternativas:
        blocos.append(LayoutBlock(tokens=[]))  # linha em branco antes das alternativas
        for alt in questao.alternativas:
            tokens_alt = [Token(text=f"{alt.letra}) "), *tokenize_rich_text(alt.texto)]
            blocos.extend(_blocos_de(tokens_alt, keep_together=True))

    return blocos


def _e_imagem_de_formula(image_id: str | None) -> bool:
    """`EQUACAO_*` (equação $$...$$ renderizada, ver latex_render.py) e
    `FORMULA_*` (recorte de página original usado como fallback fiel quando a
    transcrição LaTeX da IA não passou na validação, ver
    app.formula_resolve/pdf_extract.py) são ambas representações de fórmula —
    tratadas de forma diferente das imagens de CONTEÚDO (`IMAGEM_NN`:
    gráfico/tabela/foto real do PDF, que continua sempre ocupando um slide
    inteiro): pequenas, com altura natural conhecida, sem forçar página
    exclusiva (ver `pptx_layout.LayoutBlock.altura_natural_pt`)."""
    return (image_id or "").startswith(("EQUACAO_", "FORMULA_"))


def _altura_natural_formula_pt(
    image_id: str, imagens: dict[str, bytes], layout_config: TextLayoutConfig
) -> float | None:
    """Altura (pt) que a imagem `image_id` vai ocupar quando desenhada,
    respeitando os mesmos limites de `_preencher_slide_imagem`
    (EQUACAO_MAX_LARGURA/ALTURA, nunca mais largo que a caixa de corpo) —
    usada em `paginar_blocos` para decidir se a fórmula cabe junto com o
    texto ao redor na mesma página. None se a imagem não existir (o
    placeholder "não encontrada" de `_preencher_slide_imagem` cobre esse
    caso na hora de desenhar; aqui, tratamos como altura zero/ignorável)."""
    dados = imagens.get(image_id)
    if dados is None:
        return None
    largura_nativa, altura_nativa = dimensoes_nativas_imagem(dados)
    max_largura_pt = min(layout_config.largura_util_pt, EQUACAO_MAX_LARGURA.pt)
    max_altura_pt = EQUACAO_MAX_ALTURA.pt
    _largura_pt, altura_pt = dimensoes_proporcionais(largura_nativa, altura_nativa, max_largura_pt, max_altura_pt)
    return altura_pt


def _resolver_equacoes_dos_blocos(
    blocos: list[LayoutBlock], imagens: dict[str, bytes], layout_config: TextLayoutConfig
) -> list[LayoutBlock]:
    """Resolve tokens `is_equation` (ver `latex_render.resolver_tokens_equacao`)
    dentro de cada bloco e, se algum virar imagem (mesmo no meio de um
    parágrafo), isola-a num bloco próprio. Imagens de CONTEÚDO (gráfico/
    tabela/foto) continuam forçando página dedicada; imagens de FÓRMULA
    (`_e_imagem_de_formula`) recebem `altura_natural_pt` e passam a fluir com
    o texto ao redor (ver `pptx_layout.paginar_blocos`)."""
    resultado: list[LayoutBlock] = []
    for bloco in blocos:
        if bloco.is_image:
            if bloco.altura_natural_pt is None and _e_imagem_de_formula(bloco.image_id):
                altura = _altura_natural_formula_pt(bloco.image_id, imagens, layout_config)
                resultado.append(replace(bloco, altura_natural_pt=altura))
            else:
                resultado.append(bloco)
            continue

        tokens_resolvidos = resolver_tokens_equacao(bloco.tokens, imagens)
        if not any(t.is_image for t in tokens_resolvidos):
            resultado.append(replace(bloco, tokens=tokens_resolvidos))
            continue

        atual: list[Token] = []
        for token in tokens_resolvidos:
            if token.is_image:
                if atual:
                    resultado.append(replace(bloco, tokens=atual))
                    atual = []
                altura = _altura_natural_formula_pt(token.image_id, imagens, layout_config)
                resultado.append(LayoutBlock(is_image=True, image_id=token.image_id, altura_natural_pt=altura))
            else:
                atual.append(token)
        if atual:
            resultado.append(replace(bloco, tokens=atual))

    return resultado


def _preencher_slide_imagem(slide, image_id: str, imagens: dict[str, bytes]) -> None:
    corpo_shape = encontrar_shape_por_nome(slide, "ph_corpo_questao")
    dados = imagens.get(image_id)

    if corpo_shape is None:
        return

    if dados is None:
        escrever_texto_formatado(corpo_shape.text_frame, [Token(text=f"[{image_id} não encontrada]")])
        return

    left, top, width, height = corpo_shape.left, corpo_shape.top, corpo_shape.width, corpo_shape.height
    corpo_shape.text_frame.clear()

    if _e_imagem_de_formula(image_id):
        # fórmula (renderizada ou recorte-fallback) é pequena — usar a área
        # toda do corpo faria uma imagem pequena esticar até ficar enorme e
        # borrada
        largura = min(width, EQUACAO_MAX_LARGURA)
        altura = min(height, EQUACAO_MAX_ALTURA)
        left = left + (width - largura) // 2
        top = top + (height - altura) // 2
        width, height = largura, altura

    adicionar_imagem_slide(slide, dados, left, top, width, height)


def _aplicar_fonte_corpo(text_frame) -> None:
    """Fonte fixa e explícita em TODOS os runs — a mesma família e tamanho
    usados pela medição em pptx_layout.py (FONTE_CORPO/TAMANHO_FONTE_CORPO_PT),
    nunca deixado para herdar do layout (que declara Arial 14pt, ver
    scripts/build_templates.py). Sem "encolher texto até caber": a paginação
    (pptx_layout.paginar_blocos) já garante que o conteúdo cabe na caixa,
    então o corpo não depende mais de autoajuste do PowerPoint."""
    for paragrafo in text_frame.paragraphs:
        for run in paragrafo.runs:
            run.font.size = TAMANHO_FONTE_CORPO
            run.font.name = FONTE_CORPO


def _preencher_slide_texto(slide, pagina: list[Token]) -> None:
    corpo_shape = encontrar_shape_por_nome(slide, "ph_corpo_questao")
    if corpo_shape is None:
        return

    escrever_texto_formatado(corpo_shape.text_frame, pagina)
    _aplicar_fonte_corpo(corpo_shape.text_frame)
    corpo_shape.text_frame.word_wrap = True


def _dividir_pagina_em_segmentos(pagina: list[Token]) -> list[tuple[str, list[Token] | str]]:
    """Divide uma página "mista" (texto + fórmula fluindo juntos, ver
    `pptx_layout.paginar_blocos`) em segmentos ordenados `("texto", tokens)` /
    `("imagem", image_id)`, nos limites de cada Token `is_image`."""
    segmentos: list[tuple[str, list[Token] | str]] = []
    atual: list[Token] = []
    for token in pagina:
        if token.is_image:
            if atual:
                segmentos.append(("texto", atual))
                atual = []
            segmentos.append(("imagem", token.image_id))
        else:
            atual.append(token)
    if atual:
        segmentos.append(("texto", atual))
    return segmentos


def _altura_texto_pt(tokens: list[Token], config_layout: TextLayoutConfig) -> float:
    """Altura real (pt) que `tokens` vai ocupar — soma as linhas simuladas de
    CADA parágrafo individualmente (`_agrupar_em_paragrafos` +
    `_simular_linhas`), em vez de simular `tokens` inteiro de uma vez só: um
    segmento de página mista pode conter vários parágrafos originais
    concatenados (com separadores `new_paragraph=True` entre eles, inseridos
    por `pptx_layout._anexar`), que `_simular_linhas` sozinha não trata como
    quebra de linha (ela espera receber os tokens de UM parágrafo por vez)."""
    return sum(len(_simular_linhas(unidade, config_layout)) for unidade in _agrupar_em_paragrafos(tokens)) * (
        config_layout.altura_linha_pt
    )


def _preencher_slide_misto(
    slide, pagina: list[Token], imagens: dict[str, bytes], config_layout: TextLayoutConfig
) -> None:
    """Preenche um slide cuja página contém TEXTO e uma ou mais imagens de
    FÓRMULA fluindo juntos (ver `pptx_layout.LayoutBlock.altura_natural_pt`)
    — diferente de `_preencher_slide_texto`/`_preencher_slide_imagem`
    (conteúdo único), aqui a página é dividida em segmentos empilhados
    verticalmente: o placeholder `ph_corpo_questao` original é reaproveitado
    (redimensionado) para o primeiro segmento, e cada segmento seguinte
    ganha sua própria shape (textbox ou imagem) posicionada logo abaixo do
    anterior, usando a MESMA altura real calculada durante a paginação —
    nunca a altura "encolhe até caber" do PowerPoint."""
    corpo_shape = encontrar_shape_por_nome(slide, "ph_corpo_questao")
    if corpo_shape is None:
        return

    left, largura = corpo_shape.left, corpo_shape.width
    margem_esquerda = corpo_shape.text_frame.margin_left
    margem_direita = corpo_shape.text_frame.margin_right
    margem_superior = corpo_shape.text_frame.margin_top
    margem_inferior = corpo_shape.text_frame.margin_bottom
    top_atual = corpo_shape.top
    primeiro_texto = True

    for tipo, conteudo in _dividir_pagina_em_segmentos(pagina):
        if tipo == "imagem":
            dados = imagens.get(conteudo)
            if dados is None:
                continue
            largura_max = min(largura, EQUACAO_MAX_LARGURA)
            imagem_shape = adicionar_imagem_slide(slide, dados, left, top_atual, largura_max, EQUACAO_MAX_ALTURA)
            top_atual = top_atual + imagem_shape.height
            continue

        altura_conteudo_pt = _altura_texto_pt(conteudo, config_layout) / MARGEM_SEGURANCA_VERTICAL
        altura_shape = Pt(altura_conteudo_pt) + margem_superior + margem_inferior

        if primeiro_texto:
            shape_texto = corpo_shape
            primeiro_texto = False
        else:
            shape_texto = slide.shapes.add_textbox(left, top_atual, largura, altura_shape)
            shape_texto.text_frame.word_wrap = True
            shape_texto.text_frame.margin_left = margem_esquerda
            shape_texto.text_frame.margin_right = margem_direita
            shape_texto.text_frame.margin_top = margem_superior
            shape_texto.text_frame.margin_bottom = margem_inferior

        shape_texto.top = top_atual
        shape_texto.width = largura
        shape_texto.height = altura_shape
        escrever_texto_formatado(shape_texto.text_frame, conteudo)
        _aplicar_fonte_corpo(shape_texto.text_frame)

        top_atual = top_atual + altura_shape


def _inserir_paginas(
    prs: Presentation,
    layout_questao,
    paginas: list[list[Token]],
    imagens: dict[str, bytes],
    cabecalho_base: str,
    cursor: int,
    layout_config: TextLayoutConfig,
) -> int:
    for p_idx, pagina in enumerate(paginas):
        slide = criar_slide_de_layout(
            prs,
            layout_questao,
            cursor,
            renomear_placeholders={IDX_HEADER_QUESTAO: "ph_header_questao", IDX_CORPO_QUESTAO: "ph_corpo_questao"},
        )
        cabecalho = cabecalho_base + (" (continuação)" if p_idx > 0 else "")

        header_shape = encontrar_shape_por_nome(slide, "ph_header_questao")
        if header_shape is not None:
            escrever_texto_formatado(header_shape.text_frame, tokenize_rich_text(cabecalho))

        imagens_na_pagina = [t for t in pagina if t.is_image]
        eh_pagina_de_imagem = len(pagina) == 1 and pagina[0].is_image
        eh_pagina_mista = len(imagens_na_pagina) >= 1 and not eh_pagina_de_imagem

        if eh_pagina_de_imagem:
            _preencher_slide_imagem(slide, pagina[0].image_id, imagens)
        elif eh_pagina_mista:
            _preencher_slide_misto(slide, pagina, imagens, layout_config)
        else:
            _preencher_slide_texto(slide, pagina)

        cursor += 1

    return cursor


def _inserir_slide_questao(
    prs: Presentation,
    layout_questao,
    disciplina: DisciplinaStats,
    questao: Questao,
    imagens: dict[str, bytes],
    layout_config: TextLayoutConfig,
    cursor: int,
) -> int:
    blocos_enunciado = _resolver_equacoes_dos_blocos(_construir_blocos_enunciado(questao), imagens, layout_config)
    paginas_enunciado = paginar_blocos(blocos_enunciado, layout_config)

    cabecalho_base = f"{disciplina.nome} | Questão {questao.numero} ▪ {questao.banca}"
    if questao.ano:
        cabecalho_base += f"/{questao.ano}"

    return _inserir_paginas(prs, layout_questao, paginas_enunciado, imagens, cabecalho_base, cursor, layout_config)


def gerar_pptx(
    extraction: ExtractionResult,
    template_path: str,
    saida_path: str,
    imagens: dict[str, bytes] | None = None,
) -> str:
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template não encontrado: {template_path}")

    imagens = imagens or {}
    prs = Presentation(template_path)
    disciplinas = build_disciplinas_stats(extraction)

    cursor, stencils_obsoletos = _localizar_slides_pre_existentes(prs)
    _substituir_tokens_no_slide(
        prs.slides[cursor], {"{{NOME_CONCURSO}}": extraction.nome_concurso or "Concurso"}
    )

    layout_capa = encontrar_layout_por_nome(prs, LAYOUT_CAPA)
    layout_contracapa = encontrar_layout_por_nome(prs, LAYOUT_CONTRACAPA)
    layout_divisor_disciplina = encontrar_layout_por_nome(prs, LAYOUT_DIVISOR_DISCIPLINA)
    layout_divisor_concurso_banca = encontrar_layout_por_nome(prs, LAYOUT_DIVISOR_CONCURSO_BANCA)
    layout_questao = encontrar_layout_por_nome(prs, LAYOUT_QUESTAO)

    # geometria REAL do corpo da questão (não números repetidos no código) —
    # lida diretamente do placeholder do layout (idêntica em todo slide
    # criado a partir dele).
    corpo_layout_placeholder = next(
        ph for ph in layout_questao.placeholders if ph.placeholder_format.idx == IDX_CORPO_QUESTAO
    )
    layout_config = TextLayoutConfig.from_shape(corpo_layout_placeholder)

    cursor = _inserir_slide_estatico(prs, layout_capa, cursor)
    cursor = _inserir_slide_estatico(prs, layout_contracapa, cursor)

    for disciplina in disciplinas:
        cursor = _inserir_divisor_disciplina(prs, layout_divisor_disciplina, disciplina, cursor)
        for grupo in _agrupar_por_concurso_banca(disciplina.questoes):
            cursor = _inserir_divisor_concurso_banca(prs, layout_divisor_concurso_banca, grupo[0], cursor)
            for questao in grupo:
                cursor = _inserir_slide_questao(
                    prs, layout_questao, disciplina, questao, imagens, layout_config, cursor
                )

    # só agora, DEPOIS de todas as inserções (ver docstring de
    # `_localizar_slides_pre_existentes` para o motivo de nunca remover antes).
    for indice in sorted(stencils_obsoletos, reverse=True):
        remover_slide(prs, indice)

    os.makedirs(os.path.dirname(saida_path), exist_ok=True)
    prs.save(saida_path)
    return saida_path

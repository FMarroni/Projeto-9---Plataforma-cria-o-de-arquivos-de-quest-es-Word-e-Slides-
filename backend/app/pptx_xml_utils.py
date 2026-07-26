"""Manipulação de OOXML de baixo nível para python-pptx, portada/adaptada de
`Projeto 1 - pptx_builder.py` e `Projeto 2 - Código.txt` (o hack de sobrescrito/
subscrito e a paginação/duplicação de slide não têm equivalente na API pública
do python-pptx)."""

import copy
import io

import fitz
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from app.formatting import Token

_EMU_POR_PIXEL_96DPI = 914400 / 96

# Mesma cor de marca usada nos DOCX (ver app/docx_common.py COR_MARCA) e no
# template PPTX corporativo (scripts/build_templates.py COR_MARCA_PPTX) —
# aplicada aqui a runs individuais marcados com `Token.cor_marca=True`.
COR_MARCA = RGBColor(0x42, 0x31, 0xA4)


def substituir_texto_mantendo_formatacao(text_frame, busca: str, substituto: str) -> None:
    """Porte verbatim do helper legado: só troca o texto dentro de runs que já
    contêm `busca` por inteiro, preservando a formatação desse run."""
    for paragraph in text_frame.paragraphs:
        if busca in paragraph.text:
            for run in paragraph.runs:
                if busca in run.text:
                    run.text = run.text.replace(busca, substituto)


def aplicar_token_run(run, token: Token) -> None:
    """Aplica um Token (de formatting.tokenize_rich_text) a um run de pptx,
    usando o hack de XML `baseline` para sobrescrito/subscrito (30000/-25000),
    exatamente como no legado."""
    run.text = token.text
    if token.bold:
        run.font.bold = True
    if token.italic:
        run.font.italic = True
    if token.underline:
        run.font.underline = True
    if token.cor_marca:
        run.font.color.rgb = COR_MARCA
    if token.superscript or token.subscript:
        rPr = run._r.get_or_add_rPr()
        rPr.set('baseline', '30000' if token.superscript else '-25000')


def escrever_texto_formatado(text_frame, tokens: list[Token]) -> None:
    """Escreve uma lista de Tokens já tokenizada (`formatting.tokenize_rich_text`)
    num text_frame de pptx, quebrando parágrafo em cada Token com
    `new_paragraph=True`. Tokens de imagem (`is_image=True`) são ignorados aqui
    — imagens de página inteira são tratadas separadamente por
    `adicionar_imagem_slide`, já que python-pptx não suporta imagem inline
    dentro de um run de texto."""
    text_frame.clear()
    paragrafo = text_frame.paragraphs[0]
    primeiro = True
    for token in tokens:
        if token.is_image:
            continue
        if token.new_paragraph:
            paragrafo = text_frame.add_paragraph()
            primeiro = False
            continue
        if not token.text and not primeiro:
            continue
        run = paragrafo.add_run()
        aplicar_token_run(run, token)


def dimensoes_proporcionais(
    largura_nativa: float, altura_nativa: float, max_largura: float, max_altura: float
) -> tuple[float, float]:
    """Calcula (largura, altura) escaladas para caber dentro de uma caixa
    (max_largura, max_altura), preservando a proporção original — agnóstico
    de unidade (EMU, pt, px...), desde que todos os 4 argumentos usem a
    MESMA unidade. Reaproveitado por `adicionar_imagem_slide` (posicionamento
    em EMU) e por `pptx_gen.py` (altura em pt, para a paginação saber de
    antemão quanto espaço uma fórmula renderizada vai ocupar)."""
    escala = min(max_largura / largura_nativa, max_altura / altura_nativa)
    return largura_nativa * escala, altura_nativa * escala


def dimensoes_nativas_imagem(dados: bytes) -> tuple[int, int]:
    """(largura, altura) nativas de `dados` (bytes de imagem), em pixels."""
    pix = fitz.Pixmap(dados)
    return pix.width, pix.height


def adicionar_imagem_slide(slide, dados: bytes, left: int, top: int, max_width: int, max_height: int):
    """Insere `dados` (bytes de imagem) como uma shape de figura no slide,
    redimensionada para caber dentro da caixa (left, top, max_width, max_height)
    — todos os valores em EMU (ex.: `shape.left`, `shape.width` de outra shape
    do próprio slide) — preservando a proporção original e centralizando
    dentro da caixa. Usa PyMuPDF só para ler as dimensões nativas da imagem."""
    largura_nativa_px, altura_nativa_px = dimensoes_nativas_imagem(dados)
    largura_nativa_emu = largura_nativa_px * _EMU_POR_PIXEL_96DPI
    altura_nativa_emu = altura_nativa_px * _EMU_POR_PIXEL_96DPI
    largura_final, altura_final = dimensoes_proporcionais(
        largura_nativa_emu, altura_nativa_emu, max_width, max_height
    )
    largura_final, altura_final = int(largura_final), int(altura_final)
    left_centralizado = int(left + (max_width - largura_final) / 2)
    top_centralizado = int(top + (max_height - altura_final) / 2)

    return slide.shapes.add_picture(
        io.BytesIO(dados), left_centralizado, top_centralizado, width=largura_final, height=altura_final
    )


def _reposicionar_ultimo_slide(prs, destino: int) -> None:
    """`Presentation.slides.add_slide()` sempre anexa no fim — reposiciona o
    sldId recém-criado para `destino` (mesma técnica usada tanto por
    `duplicar_slide` quanto por `criar_slide_de_layout`)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    novo_sld = slides[-1]
    xml_slides.remove(novo_sld)
    xml_slides.insert(destino, novo_sld)


def duplicar_slide(prs, source_index: int, insert_at: int | None = None) -> int:
    """Duplica o slide em `source_index`. python-pptx não tem `Slide.duplicate()`
    — isto reimplementa a técnica padrão da comunidade: deepcopy das shapes +
    reposicionamento do sldId. Por padrão insere logo após o original
    (`source_index + 1`); passe `insert_at` para controlar a posição final
    (necessário ao montar várias páginas em sequência a partir do mesmo
    slide-base sem inverter a ordem). Retorna o índice do novo slide."""
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)

    # remove shapes herdadas do layout (placeholders vazios) antes de copiar
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    for shape in source.shapes:
        novo_el = copy.deepcopy(shape._element)
        dest.shapes._spTree.insert_element_before(novo_el, 'p:extLst')

    # melhor esforço: copia relacionamentos não-triviais (ex.: imagens), se houver
    try:
        for rel in source.part.rels.values():
            if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
                continue
            dest.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
    except Exception:
        pass

    destino = insert_at if insert_at is not None else source_index + 1
    _reposicionar_ultimo_slide(prs, destino)
    return destino


def criar_slide_de_layout(prs, layout, insert_at: int, renomear_placeholders: dict[int, str] | None = None):
    """Cria um slide NOVO a partir de um layout do slide mestre (`add_slide`,
    não uma cópia de slide existente) — usado para os layouts autorados
    diretamente no PowerPoint (ver `templates/slides.pptx`), cujos
    placeholders materializam corretamente como shapes editáveis nesse
    caminho (diferente de um deepcopy de slide, que não os recria). Renomeia
    os placeholders indicados (`{idx: nome}`) para os nomes convencionais
    usados por `encontrar_shape_por_nome` no restante do módulo, e
    reposiciona o slide em `insert_at` (`add_slide` sempre anexa no fim).
    Retorna o slide criado (a referência permanece válida após reposicionar)."""
    slide = prs.slides.add_slide(layout)
    for idx, nome in (renomear_placeholders or {}).items():
        placeholder = next(p for p in slide.placeholders if p.placeholder_format.idx == idx)
        placeholder.name = nome
    _reposicionar_ultimo_slide(prs, insert_at)
    return slide


def remover_slide(prs, index: int) -> None:
    """Remove o slide em `index` (usado para descartar os slides-stencil
    pristinos após a duplicação, já que eles não fazem parte do conteúdo final)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    sld = slides[index]
    rId = sld.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(sld)


def duplicar_linha_tabela(table):
    """python-pptx não expõe `table.add_row()` publicamente — clona o `<a:tr>`
    da última linha (preservando número/estilo de células) e o anexa ao final
    da tabela. Retorna a nova linha. Nota: `table.rows` do python-pptx não
    aceita índice negativo (levanta IndexError), por isso usamos `len(...)-1`."""
    tbl = table._tbl
    linhas = tbl.findall(qn('a:tr'))
    ultima = linhas[-1]
    nova = copy.deepcopy(ultima)
    tbl.append(nova)
    return table.rows[len(table.rows) - 1]


def remover_linha_tabela(table, row_idx: int) -> None:
    tbl = table._tbl
    linhas = tbl.findall(qn('a:tr'))
    tbl.remove(linhas[row_idx])


def encontrar_shape_por_nome(slide, nome: str):
    for shape in slide.shapes:
        if shape.name == nome:
            return shape
    return None


def encontrar_layout_por_nome(prs, nome: str):
    """Localiza um layout do slide mestre pelo nome dado a ele no PowerPoint
    (Modo de Exibição Mestre) — ex.: 'Capa', 'Capa_Disciplina'."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == nome:
                return layout
    raise ValueError(f"Layout '{nome}' não encontrado no slide mestre do template.")


def remover_shape(shape) -> None:
    """Remove uma shape do slide (ex.: descartar do stencil um textbox/tabela
    que não deve aparecer na versão final gerada)."""
    shape._element.getparent().remove(shape._element)


def encontrar_tabela_por_nome(slide, nome: str):
    shape = encontrar_shape_por_nome(slide, nome)
    if shape is not None and shape.has_table:
        return shape.table
    return None
